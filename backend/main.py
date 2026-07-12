"""
Dockbench backend — optional, self-hosted, zero-per-call-cost API.

Every dependency here is free and open-source. There is no paid vendor
API anywhere in this file. You run this yourself (a VPS, a Docker
container, a Raspberry Pi with enough RAM) and it costs whatever your
own hosting costs — never a per-request fee.

The heavy lifting is done by battle-tested C/C++ engines driven from
Python — "C for the hot path, Python for the orchestration":

  Tesseract (C++)        -> OCR
  MuPDF (C, via PyMuPDF) -> PDF rendering, table detection
  qpdf (C++, via pikepdf)-> PDF assembly for searchable OCR output
  Ghostscript (C)        -> PDF/A conversion, deep PDF compression
  LibreOffice (C++)      -> Office <-> PDF conversion
  onnxruntime (C++)      -> background-removal inference
  native/imgproc.c       -> our own scan-cleanup kernel (see native_ops.py)

Endpoints:
  POST /ocr           -> text OR searchable PDF from an image/scanned PDF
  POST /remove-bg     -> remove image background (rembg, open-weight model)
  POST /summarize     -> summarize text using a local Ollama model
  POST /chat          -> answer a question against supplied document text
  POST /convert       -> Office <-> PDF conversions (LibreOffice/pdf2docx/...)
  POST /compress-pdf  -> deep, Ghostscript-grade PDF compression
  POST /pdfa          -> convert to ISO 19005 PDF/A-2b (archival)
  GET  /capabilities  -> which optional engines are installed
  GET  /health        -> liveness check

Privacy contract:
  - Nothing is ever logged about file contents or filenames.
  - Pure-Python endpoints process uploads entirely in memory.
  - The external engines (LibreOffice, Ghostscript) require real files,
    so those endpoints use a private per-request temporary directory
    that is deleted the moment the response is built.
  - No request data is retained after the response is sent.
"""

import io
import logging
import os
import shutil
import subprocess
import tempfile
import time
import uuid

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, Response
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded

import native_ops

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s [%(request_id)s] %(message)s",
)
logger = logging.getLogger("dockbench")

limiter = Limiter(key_func=get_remote_address, default_limits=["60/minute"])
app = FastAPI(title="Dockbench API", version="0.2.0")
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)

# Lock this down to your actual frontend origin in production.
app.add_middleware(
    CORSMiddleware,
    allow_origins=os.environ.get("ALLOWED_ORIGINS", "*").split(","),
    allow_methods=["POST", "GET"],
    allow_headers=["*"],
)


@app.middleware("http")
async def request_id_and_timing(request: Request, call_next):
    """Tags every request with a short ID for tracing in logs, and times it.
    Never logs file contents or filenames — only method, path, status, timing."""
    req_id = uuid.uuid4().hex[:8]
    start = time.time()
    response = await call_next(request)
    duration_ms = (time.time() - start) * 1000
    logging.LoggerAdapter(logger, {"request_id": req_id}).info(
        f"{request.method} {request.url.path} -> {response.status_code} ({duration_ms:.0f}ms)"
    )
    response.headers["X-Request-ID"] = req_id
    return response


MAX_FILE_MB = int(os.environ.get("MAX_FILE_MB", "50"))


def _check_size(data: bytes):
    if len(data) > MAX_FILE_MB * 1024 * 1024:
        raise HTTPException(413, f"File exceeds {MAX_FILE_MB}MB limit.")


def _which(binary: str):
    return shutil.which(binary)


def _run(cmd, timeout=180):
    """Run an external engine. Raises 422 with the engine's stderr tail on
    failure — enough to debug, never echoing document content."""
    try:
        proc = subprocess.run(
            cmd, capture_output=True, timeout=timeout, check=False
        )
    except subprocess.TimeoutExpired:
        raise HTTPException(504, "Conversion engine timed out.")
    if proc.returncode != 0:
        tail = (proc.stderr or b"")[-400:].decode("utf-8", "replace")
        raise HTTPException(422, f"Engine failed (exit {proc.returncode}): {tail}")
    return proc


@app.get("/health")
def health():
    return {"status": "ok"}


@app.get("/capabilities")
def capabilities():
    """Lets the frontend discover which optional engines this deployment
    has, so it can enable exactly the tools that will actually work."""
    return {
        "version": app.version,
        "native_c_kernel": native_ops.native_available(),
        "ocr": _which("tesseract") is not None,
        "ghostscript": _which("gs") is not None,
        "libreoffice": _which("soffice") is not None,
        "ollama_configured": bool(os.environ.get("OLLAMA_URL", "http://localhost:11434")),
        "max_file_mb": MAX_FILE_MB,
    }


# ---------------------------------------------------------------------
# OCR — Tesseract (C++), fronted by our C scan-cleanup kernel.
#   output=text -> {"text": ...}
#   output=pdf  -> a searchable PDF: the original page images with an
#                  invisible, selectable text layer (what Adobe calls
#                  "Recognize Text"). Assembled with pikepdf (qpdf, C++).
# ---------------------------------------------------------------------
@app.post("/ocr")
@limiter.limit("10/minute")
async def ocr(
    request: Request,
    file: UploadFile = File(...),
    language: str = Form("eng"),
    output: str = Form("text"),
    enhance: bool = Form(True),
):
    import pytesseract
    from PIL import Image

    if output not in ("text", "pdf"):
        raise HTTPException(422, "output must be 'text' or 'pdf'")
    # Tesseract language codes are alphanumeric plus underscore/plus
    # (e.g. "eng", "deu", "eng+fra"); reject anything else before it
    # reaches a command line.
    if not language or not all(c.isalnum() or c in "+_" for c in language):
        raise HTTPException(422, "Invalid language code.")

    data = await file.read()
    _check_size(data)
    filename = (file.filename or "").lower()

    def page_images():
        if filename.endswith(".pdf"):
            import fitz  # PyMuPDF — MuPDF's C core

            doc = fitz.open(stream=data, filetype="pdf")
            try:
                for page in doc:
                    pix = page.get_pixmap(dpi=200)
                    yield Image.open(io.BytesIO(pix.tobytes("png")))
            finally:
                doc.close()
        else:
            yield Image.open(io.BytesIO(data))

    try:
        if output == "text":
            pages_text = []
            for img in page_images():
                if enhance:
                    img = native_ops.enhance_for_ocr(img)
                pages_text.append(pytesseract.image_to_string(img, lang=language))
            return JSONResponse({"text": "\n\n--- page break ---\n\n".join(pages_text)})

        # output == "pdf": one searchable-PDF page per input page, merged.
        import pikepdf

        merged = pikepdf.Pdf.new()
        for img in page_images():
            ocr_input = native_ops.enhance_for_ocr(img) if enhance else img
            page_pdf = pytesseract.image_to_pdf_or_hocr(
                ocr_input, lang=language, extension="pdf"
            )
            with pikepdf.open(io.BytesIO(page_pdf)) as part:
                merged.pages.extend(part.pages)
        out = io.BytesIO()
        merged.save(out)
        merged.close()
        return Response(content=out.getvalue(), media_type="application/pdf")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(422, f"Could not process file: {e}")
    finally:
        del data


# ---------------------------------------------------------------------
# Background removal — open-source rembg (Apache-2.0), self-hosted.
# ---------------------------------------------------------------------
@app.post("/remove-bg")
@limiter.limit("10/minute")
async def remove_bg(request: Request, file: UploadFile = File(...)):
    from rembg import remove

    data = await file.read()
    _check_size(data)

    try:
        output = remove(data)
    except Exception as e:
        raise HTTPException(422, f"Could not process image: {e}")
    finally:
        del data

    return Response(content=output, media_type="image/png")


# ---------------------------------------------------------------------
# Office <-> PDF conversion.
#   docx/pptx/xlsx/odt/... -> pdf : LibreOffice headless (full fidelity)
#   pdf -> docx                   : pdf2docx (layout-aware, PyMuPDF core)
#   pdf -> pptx                   : page renders on real slides
#   pdf -> xlsx                   : MuPDF table detection -> openpyxl
# ---------------------------------------------------------------------
OFFICE_EXTS = {".doc", ".docx", ".odt", ".rtf", ".ppt", ".pptx", ".odp",
               ".xls", ".xlsx", ".ods", ".csv", ".txt", ".html"}


@app.post("/convert")
@limiter.limit("10/minute")
async def convert(
    request: Request,
    file: UploadFile = File(...),
    target: str = Form(...),
):
    data = await file.read()
    _check_size(data)
    filename = file.filename or "input"
    ext = os.path.splitext(filename)[1].lower()
    target = target.lower().lstrip(".")

    if target == "pdf" and ext in OFFICE_EXTS:
        return _office_to_pdf(data, ext)
    if ext == ".pdf" and target == "docx":
        return _pdf_to_docx(data)
    if ext == ".pdf" and target == "pptx":
        return _pdf_to_pptx(data)
    if ext == ".pdf" and target == "xlsx":
        return _pdf_to_xlsx(data)
    raise HTTPException(
        422,
        f"Unsupported conversion: {ext or 'unknown'} -> {target}. "
        f"Supported: Office documents -> pdf, and pdf -> docx/pptx/xlsx.",
    )


def _office_to_pdf(data: bytes, ext: str) -> Response:
    if not _which("soffice"):
        raise HTTPException(501, "LibreOffice is not installed on this server (see backend/Dockerfile).")
    with tempfile.TemporaryDirectory(prefix="dockbench-") as tmp:
        src = os.path.join(tmp, "input" + ext)
        with open(src, "wb") as f:
            f.write(data)
        # A private profile dir keeps parallel conversions from fighting
        # over LibreOffice's lock files.
        profile = os.path.join(tmp, "profile")
        _run([
            "soffice", "--headless", "--norestore",
            f"-env:UserInstallation=file://{profile}",
            "--convert-to", "pdf", "--outdir", tmp, src,
        ], timeout=180)
        out_path = os.path.join(tmp, "input.pdf")
        if not os.path.exists(out_path):
            raise HTTPException(422, "Conversion produced no output — is the file a valid document?")
        with open(out_path, "rb") as f:
            return Response(content=f.read(), media_type="application/pdf")


def _pdf_to_docx(data: bytes) -> Response:
    from pdf2docx import Converter

    with tempfile.TemporaryDirectory(prefix="dockbench-") as tmp:
        src = os.path.join(tmp, "input.pdf")
        dst = os.path.join(tmp, "output.docx")
        with open(src, "wb") as f:
            f.write(data)
        try:
            cv = Converter(src)
            cv.convert(dst)
            cv.close()
        except Exception as e:
            raise HTTPException(422, f"Could not convert this PDF to Word: {e}")
        with open(dst, "rb") as f:
            return Response(
                content=f.read(),
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )


def _pdf_to_pptx(data: bytes) -> Response:
    import fitz
    from pptx import Presentation
    from pptx.util import Emu

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as e:
        raise HTTPException(422, f"Could not open PDF: {e}")
    prs = Presentation()
    blank = prs.slide_layouts[6]
    EMU_PER_POINT = 12700
    try:
        first = doc[0].rect
        prs.slide_width = Emu(int(first.width * EMU_PER_POINT))
        prs.slide_height = Emu(int(first.height * EMU_PER_POINT))
        for page in doc:
            pix = page.get_pixmap(dpi=150)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                io.BytesIO(pix.tobytes("png")), 0, 0,
                width=prs.slide_width, height=prs.slide_height,
            )
    finally:
        doc.close()
    out = io.BytesIO()
    prs.save(out)
    return Response(
        content=out.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


def _pdf_to_xlsx(data: bytes) -> Response:
    import fitz
    from openpyxl import Workbook

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as e:
        raise HTTPException(422, f"Could not open PDF: {e}")
    wb = Workbook()
    wb.remove(wb.active)
    tables_found = 0
    try:
        for pageno, page in enumerate(doc, start=1):
            for tno, table in enumerate(page.find_tables(), start=1):
                ws = wb.create_sheet(title=f"p{pageno}-table{tno}"[:31])
                for row in table.extract():
                    ws.append([cell if cell is not None else "" for cell in row])
                tables_found += 1
    finally:
        doc.close()
    if tables_found == 0:
        raise HTTPException(422, "No tables detected in this PDF.")
    out = io.BytesIO()
    wb.save(out)
    return Response(
        content=out.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


# ---------------------------------------------------------------------
# Deep PDF compression — Ghostscript (C). Far stronger than what a
# browser can do: image downsampling, font subsetting, stream rewriting.
# ---------------------------------------------------------------------
GS_PRESETS = {"extreme": "/screen", "strong": "/ebook", "balanced": "/printer"}


@app.post("/compress-pdf")
@limiter.limit("10/minute")
async def compress_pdf(
    request: Request,
    file: UploadFile = File(...),
    level: str = Form("strong"),
):
    if not _which("gs"):
        raise HTTPException(501, "Ghostscript is not installed on this server (see backend/Dockerfile).")
    if level not in GS_PRESETS:
        raise HTTPException(422, f"level must be one of {sorted(GS_PRESETS)}")
    data = await file.read()
    _check_size(data)
    with tempfile.TemporaryDirectory(prefix="dockbench-") as tmp:
        src = os.path.join(tmp, "in.pdf")
        dst = os.path.join(tmp, "out.pdf")
        with open(src, "wb") as f:
            f.write(data)
        _run([
            "gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.5",
            f"-dPDFSETTINGS={GS_PRESETS[level]}",
            "-dNOPAUSE", "-dQUIET", "-dBATCH", "-dSAFER",
            f"-sOutputFile={dst}", src,
        ])
        with open(dst, "rb") as f:
            out = f.read()
    # Ghostscript occasionally "compresses" a lean file into a bigger one;
    # never hand back a worse result than the input.
    if len(out) >= len(data):
        return Response(content=data, media_type="application/pdf",
                        headers={"X-Compression": "already-optimal"})
    return Response(content=out, media_type="application/pdf",
                    headers={"X-Compression": f"{len(data)}->{len(out)}"})


# ---------------------------------------------------------------------
# PDF/A — ISO 19005 archival conversion via Ghostscript.
# ---------------------------------------------------------------------
@app.post("/pdfa")
@limiter.limit("10/minute")
async def pdfa(request: Request, file: UploadFile = File(...)):
    if not _which("gs"):
        raise HTTPException(501, "Ghostscript is not installed on this server (see backend/Dockerfile).")
    data = await file.read()
    _check_size(data)
    with tempfile.TemporaryDirectory(prefix="dockbench-") as tmp:
        src = os.path.join(tmp, "in.pdf")
        dst = os.path.join(tmp, "out.pdf")
        with open(src, "wb") as f:
            f.write(data)
        _run([
            "gs", "-dPDFA=2", "-dBATCH", "-dNOPAUSE", "-dQUIET", "-dSAFER",
            "-sColorConversionStrategy=UseDeviceIndependentColor",
            "-sDEVICE=pdfwrite", "-dPDFACompatibilityPolicy=1",
            f"-sOutputFile={dst}", src,
        ])
        with open(dst, "rb") as f:
            return Response(content=f.read(), media_type="application/pdf")


# ---------------------------------------------------------------------
# Summarize / Chat — local LLM via Ollama (self-hosted, no API key, no
# per-token cost). Requires `ollama serve` running with a pulled model,
# e.g. `ollama pull llama3.2`.
# ---------------------------------------------------------------------
async def _ollama_generate(prompt: str) -> str:
    import httpx

    ollama_url = os.environ.get("OLLAMA_URL", "http://localhost:11434")
    model = os.environ.get("OLLAMA_MODEL", "llama3.2")
    try:
        async with httpx.AsyncClient(timeout=120) as client:
            r = await client.post(
                f"{ollama_url}/api/generate",
                json={"model": model, "prompt": prompt, "stream": False},
            )
            r.raise_for_status()
            return r.json().get("response", "").strip()
    except Exception as e:
        raise HTTPException(
            503,
            f"Local model unavailable — is Ollama running with a pulled model? ({e})",
        )


@app.post("/summarize")
@limiter.limit("20/minute")
async def summarize(request: Request, text: str = Form(...), max_words: int = Form(150)):
    max_words = max(20, min(int(max_words), 1000))
    summary = await _ollama_generate(
        f"Summarize the following text in no more than {max_words} words. "
        f"Be concise and factual.\n\n{text[:60000]}"
    )
    return JSONResponse({"summary": summary})


@app.post("/chat")
@limiter.limit("20/minute")
async def chat(request: Request, question: str = Form(...), context: str = Form(...)):
    """Chat-with-your-document. The frontend extracts the text on-device and
    sends only the question plus the relevant excerpt — never the file."""
    answer = await _ollama_generate(
        "You are a helpful assistant answering questions about a document. "
        "Use ONLY the document excerpt below. If the answer is not in the "
        "excerpt, say so plainly.\n\n"
        f"DOCUMENT EXCERPT:\n{context[:60000]}\n\nQUESTION: {question[:2000]}\n\nANSWER:"
    )
    return JSONResponse({"answer": answer})
