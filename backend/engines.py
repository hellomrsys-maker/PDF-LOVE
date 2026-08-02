"""
engines — the actual document-processing logic, framework-free.

Every heavy operation lives here as a pure function taking bytes and
options and returning `(payload, media_type)`. Both callers share them:

  main.py    -> synchronous API endpoints (small deployments, dev)
  worker.py  -> ARQ background jobs (production: the API stays responsive
                while OCR/LibreOffice/Ghostscript grind in worker
                containers)

Errors are raised as EngineError(status, detail) so neither caller needs
the other's framework.
"""

import io
import os
import shutil
import subprocess
import tempfile

import native_ops
import logging
from typing import List, Tuple
import pikepdf


# Defined in errors.py (a leaf module) to keep engines/premium_pdf/ai_helper
# from importing each other in a cycle; re-exported here so the existing
# `from engines import EngineError` callers keep working.
from errors import EngineError


def which(binary: str):
    return shutil.which(binary)


def _run(cmd, timeout=180):
    """Run an external engine. Raises EngineError with the engine's stderr
    tail on failure — enough to debug, never echoing document content."""
    try:
        proc = subprocess.run(cmd, capture_output=True, timeout=timeout, check=False)
    except subprocess.TimeoutExpired:
        raise EngineError(504, "Conversion engine timed out.")
    if proc.returncode != 0:
        tail = (proc.stderr or b"")[-400:].decode("utf-8", "replace")
        raise EngineError(422, f"Engine failed (exit {proc.returncode}): {tail}")
    return proc


# ---------------------------------------------------------------------
# OCR — Tesseract (C++), fronted by our C scan-cleanup kernel.
# ---------------------------------------------------------------------
def run_ocr(data: bytes, filename: str, language: str = "eng",
            output: str = "text", enhance: bool = True):
    import pytesseract
    from PIL import Image

    if output not in ("text", "pdf"):
        raise EngineError(422, "output must be 'text' or 'pdf'")
    if not language or not all(c.isalnum() or c in "+_" for c in language):
        raise EngineError(422, "Invalid language code.")
    if not which("tesseract"):
        raise EngineError(501, "Tesseract is not installed on this server (see backend/Dockerfile).")

    filename = (filename or "").lower()

    def page_images():
        if filename.endswith(".pdf"):
            import fitz  # PyMuPDF — MuPDF's C core

            doc = fitz.open(stream=data, filetype="pdf")
            try:
                for page in doc:
                    pix = page.get_pixmap(dpi=200)
                    yield Image.open(io.BytesIO(pix.tobytes("png")))
            finally:
                doc.close()
        else:
            yield Image.open(io.BytesIO(data))

    try:
        if output == "text":
            pages_text = []
            for img in page_images():
                if enhance:
                    img = native_ops.enhance_for_ocr(img)
                pages_text.append(pytesseract.image_to_string(img, lang=language))
            return "\n\n--- page break ---\n\n".join(pages_text), "text/plain"

        merged = pikepdf.Pdf.new()
        for img in page_images():
            ocr_input = native_ops.enhance_for_ocr(img) if enhance else img
            page_pdf = pytesseract.image_to_pdf_or_hocr(
                ocr_input, lang=language, extension="pdf"
            )
            with pikepdf.open(io.BytesIO(page_pdf)) as part:
                merged.pages.extend(part.pages)
        out = io.BytesIO()
        merged.save(out)
        merged.close()
        return out.getvalue(), "application/pdf"
    except EngineError:
        raise
    except Exception as e:
        raise EngineError(422, f"Could not process file: {e}")


# ---------------------------------------------------------------------
# Background removal — rembg (Apache-2.0).
# ---------------------------------------------------------------------
def run_remove_bg(data: bytes):
    from rembg import remove

    try:
        return remove(data), "image/png"
    except Exception as e:
        raise EngineError(422, f"Could not process image: {e}")


# ---------------------------------------------------------------------
# Office <-> PDF conversion.
# ---------------------------------------------------------------------
OFFICE_EXTS = {".doc", ".docx", ".odt", ".rtf", ".ppt", ".pptx", ".odp",
               ".xls", ".xlsx", ".ods", ".csv", ".txt", ".html"}

MEDIA_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


def run_convert(data: bytes, filename: str, target: str):
    ext = os.path.splitext(filename or "input")[1].lower()
    target = (target or "").lower().lstrip(".")

    if target == "pdf" and ext in OFFICE_EXTS:
        return _office_to_pdf(data, ext), MEDIA_TYPES["pdf"]
    if ext == ".pdf" and target == "docx":
        return _pdf_to_docx(data), MEDIA_TYPES["docx"]
    if ext == ".pdf" and target == "pptx":
        return _pdf_to_pptx(data), MEDIA_TYPES["pptx"]
    if ext == ".pdf" and target == "xlsx":
        return _pdf_to_xlsx(data), MEDIA_TYPES["xlsx"]
    raise EngineError(
        422,
        f"Unsupported conversion: {ext or 'unknown'} -> {target}. "
        f"Supported: Office documents -> pdf, and pdf -> docx/pptx/xlsx.",
    )


def _office_to_pdf(data: bytes, ext: str) -> bytes:
    if not which("soffice"):
        raise EngineError(501, "LibreOffice is not installed on this server (see backend/Dockerfile).")
    with tempfile.TemporaryDirectory(prefix="dockbench-") as tmp:
        src = os.path.join(tmp, "input" + ext)
        with open(src, "wb") as f:
            f.write(data)
        # A private profile dir keeps parallel conversions from fighting
        # over LibreOffice's lock files.
        profile = os.path.join(tmp, "profile")
        _run([
            "soffice", "--headless", "--norestore",
            f"-env:UserInstallation=file://{profile}",
            "--convert-to", "pdf", "--outdir", tmp, src,
        ], timeout=180)
        out_path = os.path.join(tmp, "input.pdf")
        if not os.path.exists(out_path):
            raise EngineError(422, "Conversion produced no output — is the file a valid document?")
        with open(out_path, "rb") as f:
            return f.read()


def _pdf_to_docx(data: bytes) -> bytes:
    from pdf2docx import Converter

    with tempfile.TemporaryDirectory(prefix="dockbench-") as tmp:
        src = os.path.join(tmp, "input.pdf")
        dst = os.path.join(tmp, "output.docx")
        with open(src, "wb") as f:
            f.write(data)
        try:
            cv = Converter(src)
            cv.convert(dst)
            cv.close()
        except Exception as e:
            raise EngineError(422, f"Could not convert this PDF to Word: {e}")
        with open(dst, "rb") as f:
            return f.read()


def _pdf_to_pptx(data: bytes) -> bytes:
    import fitz
    from pptx import Presentation
    from pptx.util import Emu

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as e:
        raise EngineError(422, f"Could not open PDF: {e}")
    prs = Presentation()
    blank = prs.slide_layouts[6]
    EMU_PER_POINT = 12700
    try:
        first = doc[0].rect
        prs.slide_width = Emu(int(first.width * EMU_PER_POINT))
        prs.slide_height = Emu(int(first.height * EMU_PER_POINT))
        for page in doc:
            pix = page.get_pixmap(dpi=150)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                io.BytesIO(pix.tobytes("png")), 0, 0,
                width=prs.slide_width, height=prs.slide_height,
            )
    finally:
        doc.close()
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _pdf_to_xlsx(data: bytes) -> bytes:
    import fitz
    from openpyxl import Workbook

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as e:
        raise EngineError(422, f"Could not open PDF: {e}")
    wb = Workbook()
    wb.remove(wb.active)
    tables_found = 0
    try:
        for pageno, page in enumerate(doc, start=1):
            for tno, table in enumerate(page.find_tables(), start=1):
                ws = wb.create_sheet(title=f"p{pageno}-table{tno}"[:31])
                for row in table.extract():
                    ws.append([cell if cell is not None else "" for cell in row])
                tables_found += 1
    finally:
        doc.close()
    if tables_found == 0:
        raise EngineError(422, "No tables detected in this PDF.")
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


# ---------------------------------------------------------------------
# Deep PDF compression + PDF/A — Ghostscript (C).
# ---------------------------------------------------------------------
GS_PRESETS = {"extreme": "/screen", "strong": "/ebook", "balanced": "/printer"}


def run_compress_pdf(data: bytes, level: str = "strong"):
    # Normalize level to lower‑case for case‑insensitive usage
    level = level.lower()
    if not which("gs"):
        raise EngineError(503, "Ghostscript is not installed on this server (see backend/Dockerfile).")
    if level not in GS_PRESETS:
        raise EngineError(422, f"level must be one of {sorted(GS_PRESETS)}")
    logging.info("Compress‑PDF: input=%d bytes, preset=%s", len(data), level)
    with tempfile.TemporaryDirectory(prefix="dockbench-") as tmp:
        src = os.path.join(tmp, "in.pdf")
        dst = os.path.join(tmp, "out.pdf")
        with open(src, "wb") as f:
            f.write(data)
        _run([
            "gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.5",
            f"-dPDFSETTINGS={GS_PRESETS[level]}",
            "-dNOPAUSE", "-dQUIET", "-dBATCH", "-dSAFER",
            f"-sOutputFile={dst}", src,
        ])
        try:
            with open(dst, "rb") as f:
                out = f.read()
        except FileNotFoundError:
            raise EngineError(422, "Ghostscript failed to produce an output PDF.")
    logging.info("Compress‑PDF: output=%d bytes", len(out))
    # Never hand back a larger file than the input.
    return (data if len(out) >= len(data) else out), "application/pdf"


def run_pdfa(data: bytes):
    if not which("gs"):
        raise EngineError(501, "Ghostscript is not installed on this server (see backend/Dockerfile).")
    with tempfile.TemporaryDirectory(prefix="dockbench-") as tmp:
        src = os.path.join(tmp, "in.pdf")
        dst = os.path.join(tmp, "out.pdf")
        with open(src, "wb") as f:
            f.write(data)
        _run([
            "gs", "-dPDFA=2", "-dBATCH", "-dNOPAUSE", "-dQUIET", "-dSAFER",
            "-sColorConversionStrategy=UseDeviceIndependentColor",
            "-sDEVICE=pdfwrite", "-dPDFACompatibilityPolicy=1",
            f"-sOutputFile={dst}", src,
        ])
        try:
            with open(dst, "rb") as f:
                return f.read(), "application/pdf"
        except FileNotFoundError:
            raise EngineError(422, "Ghostscript failed to produce a PDF/A output.")

# ---------------------------------------------------------------------
# Video processing – FFmpeg (C).
# ---------------------------------------------------------------------
def run_video_process(data: bytes, filename: str, codec: str = "h264", quality: str = "23"):
    """Process a video file with FFmpeg.

    * ``codec`` – video codec (e.g., ``h264`` or ``vp9``).
    * ``quality`` – CRF value (0‑51, lower = higher quality).
    """
    if not which("ffmpeg"):
        raise EngineError(503, "FFmpeg is not installed on this server (see backend/Dockerfile).")
    # Ensure quality is an integer within valid range
    try:
        crf = int(quality)
        if not (0 <= crf <= 51):
            raise ValueError()
    except ValueError:
        raise EngineError(422, f"Invalid quality value '{quality}'. Must be integer 0‑51.")
    if not codec.replace("-", "").isalnum():
        raise EngineError(422, f"Invalid codec '{codec}'.")
    # Create temporary files for input and output
    with tempfile.TemporaryDirectory(prefix="dockbench-") as tmp:
        # basename() only — the uploaded filename is attacker-controlled and
        # must never be able to escape the private temp dir.
        src = os.path.join(tmp, os.path.basename(filename) or "input")
        # Pick a real container: ffmpeg infers the muxer from the extension,
        # and "out.vp9" is not one it knows.
        dst = os.path.join(tmp, "out.webm" if codec in ("vp9", "vp8", "libvpx", "libvpx-vp9") else "out.mp4")
        # Write the input bytes to the source file
        with open(src, "wb") as f:
            f.write(data)
        # Build ffmpeg command
        cmd = [
            "ffmpeg",
            "-y",                # overwrite output without asking
            "-i", src,
            "-c:v", codec,
            "-crf", str(crf),
            "-preset", "medium",
            dst,
        ]
        _run(cmd)
        # Read the processed video
        with open(dst, "rb") as f:
            out = f.read()
    media_type = "video/mp4" if dst.endswith(".mp4") else "video/webm"
    return out, media_type


from ai_helper import run_image_enhance, run_super_resolve, run_audio_extract
from premium_pdf import run_merge_pdf, run_split_pdf, run_rotate_pdf, run_watermark_pdf
# Each takes (data, filename, **options) and returns (payload, media_type).
JOB_KINDS = {
    "ocr": lambda data, filename, **o: run_ocr(
        data, filename,
        language=o.get("language", "eng"),
        output=o.get("output", "text"),
        enhance=str(o.get("enhance", "true")).lower() != "false",
    ),
    "convert": lambda data, filename, **o: run_convert(data, filename, o.get("target", "")),
    "compress-pdf": lambda data, filename, **o: run_compress_pdf(data, o.get("level", "strong")),
    "pdfa": lambda data, filename, **o: run_pdfa(data),
    "remove-bg": lambda data, filename, **o: run_remove_bg(data),
    "video-process": lambda data, filename, **o: run_video_process(data, filename, o.get("codec", "h264"), o.get("quality", "23")),
    "merge-pdf": lambda data, filename, **o: run_merge_pdf(o.get("files", [])),
    "split-pdf": lambda data, filename, **o: run_split_pdf(data),
    "rotate-pdf": lambda data, filename, **o: run_rotate_pdf(data, o.get("angle", 0)),
    "watermark-pdf": lambda data, filename, **o: run_watermark_pdf(data, o.get("watermark", b"")),
}
