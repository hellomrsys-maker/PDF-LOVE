#!/usr/bin/env python3
"""
Generate the real test files the functional suite feeds to every tool.

    python scripts/make-fixtures.py [--out DIR]

Generated rather than committed: binaries in git go stale, bloat clones, and
nobody can tell what is actually inside them. Every fixture here is built
from a known specification, so a test can assert against what went in —
"the OCR output contains the words we rendered", not "some text came back".

Fixtures deliberately include awkward cases, because those are where the
bugs are: a corrupt PDF, an encrypted PDF, a photo with GPS EXIF, a page
with an SSN on it, a document photographed at an angle.
"""

import argparse
import io
import json
import os
import random
import shutil
import subprocess
import sys
import zipfile

random.seed(1234)          # reproducible fixtures

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEFAULT_OUT = os.path.join(ROOT, "fixtures")

# Text rendered into the OCR fixture, asserted against by the OCR tests.
OCR_PHRASE = "The quick brown fox jumps over the lazy dog"
# Deliberately fake, matching the redaction tool's SSN pattern.
SECRET_SSN = "123-45-6789"
SECRET_EMAIL = "private.person@example.com"


def need(mod):
    try:
        return __import__(mod)
    except ImportError:
        sys.exit(f"missing dependency: {mod}\n  pip install pikepdf pillow reportlab")


# ---------------------------------------------------------------- PDFs

def pdf_text(path, pages=3):
    """A normal text PDF with selectable text."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(path, pagesize=letter)
    for i in range(pages):
        c.setFont("Helvetica-Bold", 18)
        c.drawString(72, 720, f"Dockbench test document — page {i + 1}")
        c.setFont("Helvetica", 11)
        y = 680
        for line in [
            OCR_PHRASE,
            "This paragraph exists so text extraction has something to find.",
            f"Social Security Number: {SECRET_SSN}",
            f"Contact: {SECRET_EMAIL}",
            "Phone: 555-010-9999",
        ]:
            c.drawString(72, y, line)
            y -= 20
        c.showPage()
    c.save()


def pdf_table(path):
    """A PDF containing a grid, for table extraction."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    rows = [["Item", "Qty", "Price"],
            ["Widget", "2", "9.99"],
            ["Gadget", "5", "24.50"],
            ["Doohickey", "1", "3.75"]]
    c = canvas.Canvas(path, pagesize=letter)
    c.setFont("Helvetica", 11)
    x0, y0, cw, rh = 72, 700, 120, 22
    for r, row in enumerate(rows):
        for col, cell in enumerate(row):
            c.drawString(x0 + col * cw + 4, y0 - r * rh + 6, cell)
        c.line(x0, y0 - r * rh, x0 + cw * 3, y0 - r * rh)
    for col in range(4):
        c.line(x0 + col * cw, y0 + rh - 16, x0 + col * cw, y0 - len(rows) * rh)
    c.save()


def pdf_scanned(path, pages=2):
    """Pages that are images of text — nothing selectable, so OCR is the
    only way to read them. This is what a real scan looks like."""
    from PIL import Image, ImageDraw, ImageFont
    import pikepdf
    from pikepdf import Name

    pdf = pikepdf.Pdf.new()
    for i in range(pages):
        img = Image.new("RGB", (1240, 1754), (250, 249, 245))
        d = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype(
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 40)
            small = ImageFont.truetype(
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 30)
        except OSError:
            font = small = ImageFont.load_default()
        d.text((90, 120), f"Scanned page {i + 1}", fill=(20, 20, 20), font=font)
        d.text((90, 220), OCR_PHRASE, fill=(25, 25, 25), font=small)
        d.text((90, 280), "1234567890", fill=(25, 25, 25), font=small)
        # A little speckle so it is not an unrealistically clean render.
        for _ in range(300):
            x, y = random.randrange(1240), random.randrange(1754)
            d.point((x, y), fill=(200, 200, 195))

        buf = io.BytesIO()
        img.save(buf, "JPEG", quality=90)
        raw = buf.getvalue()

        page = pdf.add_blank_page(page_size=(612, 792))
        st = pikepdf.Stream(pdf, raw)
        st.Type = Name("/XObject"); st.Subtype = Name("/Image")
        st.Width, st.Height = 1240, 1754
        st.ColorSpace = Name("/DeviceRGB"); st.BitsPerComponent = 8
        st.Filter = Name("/DCTDecode")
        page.Resources = pikepdf.Dictionary(XObject=pikepdf.Dictionary(Im0=st))
        page.Contents = pikepdf.Stream(pdf, b"q 612 0 0 792 0 0 cm /Im0 Do Q")
    pdf.save(path)
    pdf.close()


def pdf_form(path):
    """A PDF with fillable AcroForm fields."""
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    c = canvas.Canvas(path, pagesize=letter)
    c.setFont("Helvetica", 12)
    c.drawString(72, 720, "Name:")
    c.acroForm.textfield(name="fullname", x=140, y=714, width=260, height=20,
                         borderStyle="inset", value="")
    c.drawString(72, 680, "Email:")
    c.acroForm.textfield(name="email", x=140, y=674, width=260, height=20,
                         borderStyle="inset", value="")
    c.drawString(72, 640, "Subscribe:")
    c.acroForm.checkbox(name="subscribe", x=140, y=634, size=18)
    c.save()


def pdf_encrypted(path, password="correct-horse"):
    import pikepdf
    src = os.path.join(os.path.dirname(path), "text.pdf")
    with pikepdf.open(src) as pdf:
        pdf.save(path, encryption=pikepdf.Encryption(
            owner=password, user=password, R=6))     # R=6 = AES-256


def pdf_corrupt(path):
    """A broken cross-reference index — the realistic damage case.

    Every object and the trailer survive; only the startxref offset is wrong,
    which is what a bad download or a buggy writer actually produces. Viewers
    and pdf-lib reject it; a real repair tool rebuilds the index by scanning
    for objects and gets all pages back. This is the case Repair PDF must fix.
    """
    import re
    src = os.path.join(os.path.dirname(path), "text.pdf")
    data = open(src, "rb").read()
    broken = re.sub(rb"startxref\s*\r?\n\s*\d+", b"startxref\n999999", data)
    assert broken != data, "startxref not found — fixture would not be damaged"
    with open(path, "wb") as f:
        f.write(broken)


def pdf_truncated(path):
    """Truncated at 82%: the trailer and catalog are gone.

    Deliberately beyond recovery — qpdf 11.9 reports "unable to find trailer
    dictionary" on it, and no tool can do better, because the data is simply
    not there. Repair PDF must fail on this with an honest message rather
    than claiming success or pointing at a server.
    """
    src = os.path.join(os.path.dirname(path), "text.pdf")
    data = open(src, "rb").read()
    with open(path, "wb") as f:
        f.write(data[:int(len(data) * 0.82)])


def pdf_bookmarked(path):
    """A PDF with a real outline, so Split by Bookmark has something to split.

    Without one the tool correctly reports "no bookmarks found" and the test
    proves nothing about whether splitting works.
    """
    import pikepdf
    src = os.path.join(os.path.dirname(path), "many-pages.pdf")
    with pikepdf.open(src) as pdf:
        with pdf.open_outline() as outline:
            for title, page in (("Introduction", 0), ("Chapter One", 100),
                                ("Chapter Two", 250), ("Appendix", 400)):
                if page < len(pdf.pages):
                    # A PDF destination is an array of [page, fit-mode], not
                    # a bare page object — pdf.js and pikepdf both reject the
                    # latter when reading the outline back.
                    dest = pikepdf.Array(
                        [pdf.pages[page].obj, pikepdf.Name("/Fit")])
                    outline.root.append(
                        pikepdf.OutlineItem(title, destination=dest))
        pdf.save(path)


def pdf_many_pages(path, pages=500):
    import pikepdf
    pdf = pikepdf.Pdf.new()
    for _ in range(pages):
        pdf.add_blank_page(page_size=(612, 792))
    pdf.save(path)
    pdf.close()


def pdf_mixed_orientation(path):
    import pikepdf
    pdf = pikepdf.Pdf.new()
    pdf.add_blank_page(page_size=(612, 792))        # portrait
    pdf.add_blank_page(page_size=(792, 612))        # landscape
    pdf.save(path)
    pdf.close()


# -------------------------------------------------------------- images

def images(out):
    from PIL import Image, ImageDraw

    base = Image.new("RGB", (800, 600), (240, 236, 228))
    d = ImageDraw.Draw(base)
    d.ellipse((250, 150, 550, 450), fill=(193, 104, 58))
    d.rectangle((60, 60, 200, 200), fill=(74, 120, 86))
    d.text((60, 540), "Dockbench fixture image", fill=(30, 30, 30))

    base.save(os.path.join(out, "photo.jpg"), quality=92)
    base.save(os.path.join(out, "photo.png"))
    base.save(os.path.join(out, "photo.webp"))
    base.save(os.path.join(out, "photo.bmp"))
    base.convert("P", palette=Image.ADAPTIVE).save(os.path.join(out, "photo.gif"))

    # Animated GIF — some tools take only the first frame; that is worth
    # exercising rather than assuming.
    frames = []
    for i in range(6):
        f = Image.new("RGB", (200, 200), (250, 250, 245))
        ImageDraw.Draw(f).ellipse((20 + i * 10, 20, 120 + i * 10, 120),
                                  fill=(193, 104, 58))
        frames.append(f.convert("P", palette=Image.ADAPTIVE))
    frames[0].save(os.path.join(out, "animated.gif"), save_all=True,
                   append_images=frames[1:], duration=120, loop=0)

    # A large image, for resize/compress assertions.
    Image.new("RGB", (3000, 2000), (200, 210, 220)).save(
        os.path.join(out, "large.jpg"), quality=95)

    # A photo of a document at an angle — the scanner's real input.
    doc = Image.new("RGB", (960, 720), (59, 58, 54))
    dd = ImageDraw.Draw(doc)
    dd.polygon([(150, 150), (800, 95), (830, 600), (190, 660)], fill=(243, 241, 234))
    for y in range(190, 580, 26):
        dd.rectangle((190, y, 190 + random.randrange(300, 560), y + 8),
                     fill=(45, 45, 45))
    doc.save(os.path.join(out, "skewed-document.jpg"), quality=92)


def image_with_exif(out):
    """A JPEG carrying GPS and camera EXIF, so the stripper can be checked
    against something that is actually there."""
    from PIL import Image
    import piexif

    path = os.path.join(out, "photo-with-exif.jpg")
    Image.open(os.path.join(out, "photo.jpg")).save(path, quality=92)

    exif = {
        "0th": {
            piexif.ImageIFD.Make: b"Dockbench",
            piexif.ImageIFD.Model: b"FixtureCam 1.0",
            piexif.ImageIFD.Software: b"make-fixtures.py",
        },
        "Exif": {piexif.ExifIFD.DateTimeOriginal: b"2024:01:15 10:30:00"},
        "GPS": {
            piexif.GPSIFD.GPSLatitudeRef: b"N",
            piexif.GPSIFD.GPSLatitude: ((51, 1), (30, 1), (0, 1)),
            piexif.GPSIFD.GPSLongitudeRef: b"W",
            piexif.GPSIFD.GPSLongitude: ((0, 1), (7, 1), (0, 1)),
        },
        "1st": {}, "thumbnail": None,
    }
    piexif.insert(piexif.dump(exif), path)
    return path


# ------------------------------------------------------- media & docs

def media(out):
    """Video with an audio track, plus standalone audio.

    Produced by a real browser via MediaRecorder rather than ffmpeg — see
    scripts/make-media-fixture.js for why. The upshot is that the fixture is
    encoded by the same engine the video tools decode with, which makes it a
    more faithful input than anything ffmpeg would emit.
    """
    script = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                          "make-media-fixture.js")
    env = dict(os.environ)
    # Playwright may be installed globally rather than beside this repo;
    # make-media-fixture.js honours PLAYWRIGHT_PATH for exactly that case.
    if "PLAYWRIGHT_PATH" not in env:
        for cand in ("/opt/node22/lib/node_modules/playwright",
                     "/usr/lib/node_modules/playwright"):
            if os.path.isdir(cand):
                env["PLAYWRIGHT_PATH"] = cand
                break
    try:
        subprocess.run(["node", script, out], check=True,
                       capture_output=True, text=True, timeout=300, env=env)
        return True
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired,
            FileNotFoundError) as e:
        detail = getattr(e, "stderr", "") or str(e)
        print(f"  note: media fixtures skipped ({detail.strip()[:120]})")
        return False


def documents(out):
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from pptx.util import Inches

    doc = Document()
    doc.add_heading("Dockbench fixture document", 0)
    doc.add_paragraph(OCR_PHRASE)
    doc.add_paragraph("A second paragraph, so conversions have structure to keep.")
    t = doc.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "Key"; t.cell(0, 1).text = "Value"
    t.cell(1, 0).text = "Answer"; t.cell(1, 1).text = "42"
    doc.save(os.path.join(out, "document.docx"))

    wb = Workbook(); ws = wb.active; ws.title = "Data"
    ws.append(["Item", "Qty", "Price"])
    for row in [["Widget", 2, 9.99], ["Gadget", 5, 24.50], ["Doohickey", 1, 3.75]]:
        ws.append(row)
    wb.save(os.path.join(out, "spreadsheet.xlsx"))

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Dockbench fixture deck"
    s.placeholders[1].text = OCR_PHRASE
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Second slide"
    s2.placeholders[1].text_frame.text = "Bullet one"
    prs.save(os.path.join(out, "presentation.pptx"))

    with open(os.path.join(out, "data.csv"), "w") as f:
        f.write("item,qty,price\nWidget,2,9.99\nGadget,5,24.50\nDoohickey,1,3.75\n")
    with open(os.path.join(out, "data.json"), "w") as f:
        json.dump([{"item": "Widget", "qty": 2, "price": 9.99},
                   {"item": "Gadget", "qty": 5, "price": 24.50}], f, indent=2)
    with open(os.path.join(out, "document.md"), "w") as f:
        f.write(f"# Fixture\n\n{OCR_PHRASE}\n\n- one\n- two\n\n"
                "| a | b |\n|---|---|\n| 1 | 2 |\n\n```js\nconst x = 1;\n```\n")
    with open(os.path.join(out, "page.html"), "w") as f:
        f.write(f"<h1>Fixture</h1><p>{OCR_PHRASE}</p><ul><li>one</li></ul>")
    with open(os.path.join(out, "plain.txt"), "w") as f:
        f.write(OCR_PHRASE + "\n" + "Another line.\n" * 20)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=DEFAULT_OUT)
    args = ap.parse_args()
    out = os.path.abspath(args.out)

    need("pikepdf"); need("PIL"); need("reportlab")
    os.makedirs(out, exist_ok=True)

    print("building fixtures in", out)
    pdf_text(os.path.join(out, "text.pdf"))
    pdf_text(os.path.join(out, "text-2.pdf"), pages=2)
    pdf_table(os.path.join(out, "table.pdf"))
    pdf_scanned(os.path.join(out, "scanned.pdf"))
    pdf_form(os.path.join(out, "form.pdf"))
    pdf_encrypted(os.path.join(out, "encrypted.pdf"))
    pdf_corrupt(os.path.join(out, "corrupt.pdf"))
    pdf_truncated(os.path.join(out, "truncated.pdf"))
    pdf_many_pages(os.path.join(out, "many-pages.pdf"))
    pdf_bookmarked(os.path.join(out, "bookmarked.pdf"))
    pdf_mixed_orientation(os.path.join(out, "mixed-orientation.pdf"))
    print("  PDFs done")

    images(out)
    try:
        image_with_exif(out)
        print("  images done (including EXIF/GPS)")
    except ImportError:
        print("  images done (piexif absent — no EXIF fixture)")

    has_media = media(out)
    documents(out)
    print("  documents done")

    manifest = {
        "ocr_phrase": OCR_PHRASE,
        "secret_ssn": SECRET_SSN,
        "secret_email": SECRET_EMAIL,
        "encrypted_password": "correct-horse",
        "has_media": has_media,
        "files": sorted(f for f in os.listdir(out) if f != "manifest.json"),
    }
    with open(os.path.join(out, "manifest.json"), "w") as f:
        json.dump(manifest, f, indent=2)

    total = sum(os.path.getsize(os.path.join(out, f)) for f in manifest["files"])
    print(f"\n{len(manifest['files'])} fixtures, {total / 1048576:.1f} MB")


if __name__ == "__main__":
    main()
